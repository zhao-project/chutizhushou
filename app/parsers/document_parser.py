"""文档解析器"""

from abc import ABC, abstractmethod
from typing import Optional


class DocumentParser(ABC):
    """文档解析基类"""

    @abstractmethod
    def parse(self, file_path: str) -> str:
        """解析文档，返回纯文本内容"""
        pass


class WordParser(DocumentParser):
    """Word (.docx) 解析"""

    def parse(self, file_path: str) -> str:
        from docx import Document
        doc = Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        # 提取表格内容
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)
        return "\n\n".join(paragraphs)


class PDFParser(DocumentParser):
    """PDF 解析"""

    def parse(self, file_path: str) -> str:
        import pdfplumber
        pages = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    pages.append(text.strip())
        return "\n\n".join(pages)


class PPTXParser(DocumentParser):
    """PPTX (.pptx) 解析"""

    def parse(self, file_path: str) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides.append("\n".join(texts))
        return "\n\n--- 下一页 ---\n\n".join(slides)


def parse_document(file_path: str, file_type: Optional[str] = None) -> str:
    """自动检测文件类型并解析"""
    if file_type is None:
        file_type = file_path.rsplit(".", 1)[-1].lower()

    parser_map = {
        "docx": WordParser,
        "pdf": PDFParser,
        "pptx": PPTXParser,
        "ppt": PPTXParser,
    }

    parser_class = parser_map.get(file_type)
    if parser_class is None:
        raise ValueError(f"不支持的文件类型: {file_type}，支持: {list(parser_map.keys())}")

    return parser_class().parse(file_path)
